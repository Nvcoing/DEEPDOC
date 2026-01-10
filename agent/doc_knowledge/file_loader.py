from typing import List, Tuple
import os
import shutil
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from pdf2image import convert_from_path
from PIL import Image
from paddleocr import PaddleOCR


# ================= CONFIG =================

# PaddleOCR auto multi-language
OCR = PaddleOCR(
    use_angle_cls=True,
    lang="multi",      # tự động detect ngôn ngữ
    show_log=False
)


# ================= POPPLER AUTO =================

def get_poppler_path() -> str | None:
    """
    100% tự động:
    - ENV POPPLER_PATH
    - PATH system (pdfinfo / pdftoppm)
    """
    env_path = os.getenv("POPPLER_PATH")
    if env_path and shutil.which("pdfinfo", path=env_path):
        return env_path

    for exe in ("pdfinfo", "pdftoppm"):
        exe_path = shutil.which(exe)
        if exe_path:
            return os.path.dirname(exe_path)

    return None


POPPLER_PATH = get_poppler_path()


# ================= OCR =================

def ocr_image(image: Image.Image) -> str:
    """
    PaddleOCR expects numpy array or image path
    """
    result = OCR.ocr(image, cls=True)

    texts = []
    for line in result or []:
        for box, (text, score) in line:
            if score >= 0.5 and text.strip():
                texts.append(text.strip())

    return "\n".join(texts)


# ================= LOAD FILE =================

def load_file_pages(path: str) -> List[str]:
    ext = path.split('.')[-1].lower()
    pages: List[str] = []

    # -------- PDF --------
    if ext == "pdf":
        reader = PdfReader(path)

        for page_id, page in enumerate(reader.pages):
            text = (page.extract_text() or "").strip()

            # Nếu page không có text → OCR
            if not text:
                try:
                    images = convert_from_path(
                        path,
                        dpi=250,
                        first_page=page_id + 1,
                        last_page=page_id + 1,
                        poppler_path=POPPLER_PATH
                    )
                    if images:
                        text = ocr_image(images[0])
                except Exception as e:
                    print(f"[WARN] OCR failed page {page_id + 1}: {e}")
                    text = ""

            if text:
                pages.append(text)

        return pages

    # -------- DOCX --------
    if ext == "docx":
        doc = Document(path)
        buf, cnt = [], 0

        for para in doc.paragraphs:
            if para.text.strip():
                buf.append(para.text.strip())
                cnt += 1

            if cnt >= 20:
                page = "\n".join(buf).strip()
                if page:
                    pages.append(page)
                buf, cnt = [], 0

        if buf:
            page = "\n".join(buf).strip()
            if page:
                pages.append(page)

        return pages

    # -------- PPTX --------
    if ext == "pptx":
        pres = Presentation(path)

        for slide in pres.slides:
            texts = [
                s.text.strip()
                for s in slide.shapes
                if hasattr(s, "text") and s.text.strip()
            ]
            page = "\n".join(texts).strip()
            if page:
                pages.append(page)

        return pages

    raise ValueError(f"Unsupported file format: {ext}")


# ================= CHUNKING =================

def chunk_pages_smart(pages: List[str]) -> List[Tuple[str, List[int]]]:
    """
    - Không sinh chunk rỗng
    - Không sinh chunk quá ngắn
    - Nối OCR text + text gốc mượt
    """
    all_chunks: List[Tuple[str, List[int]]] = []

    for page_id, page_text in enumerate(pages):
        if not page_text.strip():
            continue

        words = page_text.split()
        total_words = len(words)

        if total_words < 20:
            continue

        chunk_size = max(total_words // 3, 20)

        for i in range(3):
            start = i * chunk_size
            end = total_words if i == 2 else start + chunk_size

            chunk_words = words[start:end]
            if len(chunk_words) < 20:
                continue

            chunk_text = " ".join(chunk_words)
            pages_involved = [page_id]

            # Overlap sang page sau
            if i == 2 and page_id + 1 < len(pages):
                next_words = pages[page_id + 1].split()
                overlap_size = max(len(next_words) // 5, 20)
                chunk_text += " " + " ".join(next_words[:overlap_size])
                pages_involved.append(page_id + 1)

            all_chunks.append((chunk_text.strip(), pages_involved))

    return all_chunks
