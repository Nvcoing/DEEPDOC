import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def load_file_pages(path):
    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        reader = PdfReader(path)
        return [p.extract_text() or "" for p in reader.pages]

    if ext == ".docx":
        doc = Document(path)
        pages, buf, cnt = [], [], 0
        for para in doc.paragraphs:
            if para.text.strip():
                buf.append(para.text.strip())
                cnt += 1
            if cnt >= 20:
                pages.append("\n".join(buf))
                buf, cnt = [], 0
        if buf:
            pages.append("\n".join(buf))
        return pages

    if ext == ".pptx":
        pres = Presentation(path)
        pages = []
        for slide in pres.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
            pages.append("\n".join(texts))
        return pages

    raise ValueError("Unsupported format")


def chunk_page(text, size=800, overlap=200):
    words = text.split()
    chunks, i = [], 0
    while i < len(words):
        chunks.append(" ".join(words[i:i + size]))
        i += size - overlap
    return chunks


def build_chunks(path):
    pages = load_file_pages(path)

    chunks, ids, metas = [], [], []

    for p_idx, text in enumerate(pages):
        for c_idx, c in enumerate(chunk_page(text)):
            chunks.append(c)
            ids.append(f"{os.path.basename(path)}_p{p_idx}_c{c_idx}")
            metas.append({"page": p_idx})

    return chunks, ids, metas, pages
